from __future__ import annotations
"""
Advanced OCR engine — multi-format text extraction with intelligent preprocessing.

Supported formats:
  - PDF (embedded text via pdfplumber + scanned-page OCR via Tesseract)
  - Images (JPEG, PNG, TIFF, BMP, WebP) with advanced CV2 preprocessing
  - DOCX (Word documents via python-docx — full paragraph + table extraction)
  - XLSX/XLS (Excel spreadsheets via openpyxl — all sheets, all cells)
  - Plain text / CSV / JSON / XML / HTML (direct read with encoding detection)

Image preprocessing pipeline:
  1. Grayscale conversion
  2. Noise removal (fastNlMeansDenoising)
  3. Adaptive thresholding for varied lighting
  4. Morphological operations (close small gaps in text)
  5. Contrast enhancement (CLAHE)
  6. Deskew via minAreaRect
  7. Multi-pass OCR with orientation detection

Returns a list of (page_number, text, confidence) tuples.
"""



def _extract_fields_and_tables(text: str) -> dict:
    """
    Dynamically extract key-value fields and tables from OCR text.
    Returns a dict: { 'fields': {key: value, ...}, 'tables': [table1, ...] }
    """
    import re
    fields = {}
    tables = []
    lines = [l.strip() for l in text.splitlines() if l.strip()]
    # Key-value extraction (e.g. Name: John Doe)
    kv_pattern = re.compile(r"^([A-Za-z0-9 .\-_/]+)\s*[:\-–]\s*(.+)$")
    for line in lines:
        m = kv_pattern.match(line)
        if m:
            key, value = m.group(1).strip(), m.group(2).strip()
            if key and value:
                fields[key] = value

    # Table extraction: group consecutive lines with 2+ columns (split by 2+ spaces or tabs or |)
    current_table = []
    for line in lines:
        # Split by | or 2+ spaces or tab
        if '|' in line:
            cols = [c.strip() for c in line.split('|')]
        else:
            cols = re.split(r"\s{2,}|\t", line)
        if len([c for c in cols if c]) >= 2:
            current_table.append(cols)
        else:
            if current_table:
                tables.append(current_table)
                current_table = []
    if current_table:
        tables.append(current_table)
    return {'fields': fields, 'tables': tables}
"""
Advanced OCR engine — multi-format text extraction with intelligent preprocessing.

Supported formats:
  - PDF (embedded text via pdfplumber + scanned-page OCR via Tesseract)
  - Images (JPEG, PNG, TIFF, BMP, WebP) with advanced CV2 preprocessing
  - DOCX (Word documents via python-docx — full paragraph + table extraction)
  - XLSX/XLS (Excel spreadsheets via openpyxl — all sheets, all cells)
  - Plain text / CSV / JSON / XML / HTML (direct read with encoding detection)

Image preprocessing pipeline:
  1. Grayscale conversion
  2. Noise removal (fastNlMeansDenoising)
  3. Adaptive thresholding for varied lighting
  4. Morphological operations (close small gaps in text)
  5. Contrast enhancement (CLAHE)
  6. Deskew via minAreaRect
  7. Multi-pass OCR with orientation detection

Returns a list of (page_number, text, confidence) tuples.
"""


import csv
import io
import json
import logging
import os
import re
import tempfile
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path
from typing import Any

# ====================== CRITICAL: Set PaddleOCR env vars BEFORE any imports that might use them
os.environ.setdefault("PADDLE_PDX_DISABLE_MODEL_SOURCE_CHECK", "True")
os.environ.setdefault("PADDLE_DISABLE_ANALYTICS", "True")

try:
    import cv2
    import numpy as np
    _HAS_CV2 = True
except ImportError:
    cv2 = None  # type: ignore[assignment]
    np = None  # type: ignore[assignment]
    _HAS_CV2 = False

try:
    import docx as _docx
    _HAS_DOCX = True
except ImportError:
    _docx = None  # type: ignore[assignment]
    _HAS_DOCX = False

try:
    import openpyxl
    _HAS_OPENPYXL = True
except ImportError:
    openpyxl = None  # type: ignore[assignment]
    _HAS_OPENPYXL = False

# ── New format support (all optional; degrade gracefully) ───────────────────
try:
    # Registers HEIC/HEIF as a PIL plugin so Image.open("x.heic") works.
    import pillow_heif as _pillow_heif
    _pillow_heif.register_heif_opener()
    _HAS_HEIF = True
except Exception:  # pragma: no cover - optional dep
    _HAS_HEIF = False

try:
    from pptx import Presentation as _Presentation  # type: ignore
    _HAS_PPTX = True
except ImportError:  # pragma: no cover
    _Presentation = None  # type: ignore[assignment]
    _HAS_PPTX = False

try:
    from odf.opendocument import load as _odf_load  # type: ignore
    from odf import text as _odf_text, teletype as _odf_teletype, table as _odf_table  # type: ignore
    _HAS_ODF = True
except ImportError:  # pragma: no cover
    _odf_load = None  # type: ignore[assignment]
    _HAS_ODF = False

try:
    from striprtf.striprtf import rtf_to_text as _rtf_to_text  # type: ignore
    _HAS_STRIPRTF = True
except ImportError:  # pragma: no cover
    _rtf_to_text = None  # type: ignore[assignment]
    _HAS_STRIPRTF = False

try:
    import pikepdf as _pikepdf  # type: ignore
    _HAS_PIKEPDF = True
except ImportError:  # pragma: no cover
    _pikepdf = None  # type: ignore[assignment]
    _HAS_PIKEPDF = False


import pdfplumber
import pytesseract
from PIL import Image, ImageEnhance, ImageFilter
from pytesseract.pytesseract import TesseractNotFoundError

# EasyOCR lazy globals (initialized on-demand)
_HAS_EASYOCR = False
_easyocr_reader = None

# Lazy initialization helper for EasyOCR. Avoid importing/initializing the
# reader at module import time to reduce worker startup latency and memory use.
def _ensure_easyocr_reader() -> None:
    global _HAS_EASYOCR, _easyocr_reader
    if _easyocr_reader is not None:
        return
    if not getattr(settings, "easyocr_enabled", True):
        _HAS_EASYOCR = False
        _easyocr_reader = None
        return
    try:
        import easyocr as _easyocr_mod
        langs = getattr(settings, "easyocr_languages", ["en"]) or ["en"]
        _easyocr_reader = _easyocr_mod.Reader(langs)
        _HAS_EASYOCR = True
        logger.info("EasyOCR lazily initialized (langs=%s)", langs)
    except Exception:
        _HAS_EASYOCR = False
        _easyocr_reader = None
        logger.debug("EasyOCR lazy init failed", exc_info=True)

PaddleOCR = None  # type: ignore[assignment]
PaddleOCRVL = None  # type: ignore[assignment]
_HAS_PADDLE = False
_HAS_PADDLE_VL = False
_PADDLE_IMPORT_TRIED = False

from .config import settings

logger = logging.getLogger("ocr.engine")

PageResult = tuple[int, str, float | None]
_paddle_engine = None
_paddle_engine_kind = "none"
_paddle_backend_logged = False
_tesseract_checked = False
_tesseract_available = False

# Point tesseract binary at configured path
pytesseract.pytesseract.tesseract_cmd = settings.tesseract_cmd

_IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".tiff", ".tif", ".bmp", ".webp", ".gif", ".heic", ".heif"}
_PDF_EXTENSIONS = {".pdf"}
_DOCX_EXTENSIONS = {".docx", ".doc"}
_EXCEL_EXTENSIONS = {".xlsx", ".xls"}
_PPTX_EXTENSIONS = {".pptx", ".ppt"}
_ODT_EXTENSIONS = {".odt", ".ods", ".odp"}
_RTF_EXTENSIONS = {".rtf"}
# .rtf intentionally NOT in _TEXT_EXTENSIONS — it has its own dedicated extractor that
# strips RTF control words (the previous behavior leaked raw \rtf1\ansi\... markup).
_TEXT_EXTENSIONS = {".txt", ".csv", ".json", ".xml", ".html", ".htm", ".md", ".log"}

# Public list (consumed by ingress + tests): every suffix the OCR service can ingest.
SUPPORTED_EXTENSIONS: set[str] = (
    _PDF_EXTENSIONS
    | _IMAGE_EXTENSIONS
    | _DOCX_EXTENSIONS
    | _EXCEL_EXTENSIONS
    | _PPTX_EXTENSIONS
    | _ODT_EXTENSIONS
    | _RTF_EXTENSIONS
    | _TEXT_EXTENSIONS
)


def _ensure_paddle_imported() -> None:
    global PaddleOCR, PaddleOCRVL, _HAS_PADDLE, _HAS_PADDLE_VL, _PADDLE_IMPORT_TRIED
    if _PADDLE_IMPORT_TRIED:
        return
    _PADDLE_IMPORT_TRIED = True
    os.environ.setdefault("PADDLE_PDX_DISABLE_MODEL_SOURCE_CHECK", "True")
    os.environ.setdefault("PADDLE_DISABLE_ANALYTICS", "True")
    try:
        from paddleocr import PaddleOCR as _PaddleOCR
        PaddleOCR = _PaddleOCR  # type: ignore[assignment]
        _HAS_PADDLE = True
        logger.info("[OCR] PaddleOCR imported successfully")
    except Exception as e:
        logger.warning("[OCR] PaddleOCR import failed: %s", e, exc_info=True)
        _HAS_PADDLE = False

    if not settings.enable_paddle_vl:
        PaddleOCRVL = None  # type: ignore[assignment]
        _HAS_PADDLE_VL = False
        return

    try:
        from paddleocr import PaddleOCRVL as _PaddleOCRVL
        PaddleOCRVL = _PaddleOCRVL  # type: ignore[assignment]
        _HAS_PADDLE_VL = True
    except Exception:
        # Some PaddleOCR builds expose doc-parser/VL via the PaddleOCR class itself.
        # If VL-specific class is absent, allow trying VL constructor args on PaddleOCR.
        if _HAS_PADDLE:
            PaddleOCRVL = PaddleOCR  # type: ignore[assignment]
            _HAS_PADDLE_VL = True
        else:
            PaddleOCRVL = None  # type: ignore[assignment]
            _HAS_PADDLE_VL = False


def _get_paddle_engine():
    global _paddle_engine, _paddle_engine_kind, _paddle_backend_logged
    if _paddle_engine is not None:
        return _paddle_engine
    if not settings.enable_paddle_ocr:
        return None
    _ensure_paddle_imported()
    if not _paddle_backend_logged:
        logger.info(
            "OCR backend probe: enable_paddle_vl=%s has_paddle=%s has_paddle_vl=%s",
            settings.enable_paddle_vl,
            _HAS_PADDLE,
            _HAS_PADDLE_VL,
        )
        logger.info(
            "OCR backend mode: %s",
            "vl" if settings.enable_paddle_vl else "classic (VL disabled by config)",
        )
        _paddle_backend_logged = True
    if not _HAS_PADDLE and not _HAS_PADDLE_VL:
        return None

    if settings.enable_paddle_vl and not _HAS_PADDLE_VL:
        logger.info("PaddleOCR-VL requested but disabled; using classic PaddleOCR")

    if settings.enable_paddle_vl and _HAS_PADDLE_VL:
        vl_attempts = [
            {
                "lang": settings.paddle_language,
                "show_log": False,
                "use_doc_parser": settings.paddle_vl_doc_parser,
                "enable_table_merge": settings.paddle_vl_merge_cross_page_tables,
                "enable_mkldnn": True,
                "use_onnx": True,
            },
            {
                "lang": settings.paddle_language,
                "show_log": False,
                "enable_mkldnn": True,
                "use_onnx": True,
            },
            {
                "lang": settings.paddle_language,
                "enable_mkldnn": True,
                "use_onnx": True,
            },
            {},
        ]
        last_vl_error: Exception | None = None
        for kwargs in vl_attempts:
            try:
                _paddle_engine = PaddleOCRVL(**kwargs)
                _paddle_engine_kind = "vl"
                logger.info("PaddleOCR-VL initialized (doc-parser markdown mode)")
                return _paddle_engine
            except TypeError:
                continue
            except ValueError as exc:
                # Different PaddleOCR versions accept different constructor args.
                if "Unknown argument" in str(exc):
                    last_vl_error = exc
                    continue
                last_vl_error = exc
                break
            except Exception as exc:
                last_vl_error = exc
                break
        if last_vl_error is not None:
            logger.warning("PaddleOCR-VL init failed; trying classic PaddleOCR (%s)", last_vl_error)

    classic_attempts = [
        {
            "use_textline_orientation": False,
            "use_doc_orientation_classify": False,
            "use_doc_unwarping": False,
            "lang": settings.paddle_language,
            "enable_mkldnn": True,
            "use_onnx": True,
        },
        {"lang": settings.paddle_language, "enable_mkldnn": True, "use_onnx": True},
        {},
    ]
    last_classic_error: Exception | None = None
    for idx, kwargs in enumerate(classic_attempts):
        try:
            logger.debug("[OCR] PaddleOCR classic attempt %d: %s", idx + 1, kwargs)
            _paddle_engine = PaddleOCR(**kwargs)
            _paddle_engine_kind = "classic"
            logger.info("[OCR] PaddleOCR initialized successfully")
            return _paddle_engine
        except ValueError as exc:
            if "Unknown argument" in str(exc):
                logger.debug("[OCR] PaddleOCR attempt %d failed (unknown arg): %s", idx + 1, exc)
                last_classic_error = exc
                continue
            logger.debug("[OCR] PaddleOCR attempt %d failed (value error): %s", idx + 1, exc)
            last_classic_error = exc
            break
        except Exception as exc:
            logger.debug("[OCR] PaddleOCR attempt %d failed (other): %s", idx + 1, exc)
            last_classic_error = exc
            break

    if last_classic_error is not None:
        logger.warning("[OCR] All PaddleOCR init attempts failed: %s", last_classic_error)
    _paddle_engine = None
    _paddle_engine_kind = "none"
    return _paddle_engine


def _extract_markdown_from_vl_payload(payload: Any) -> str:
    if payload is None:
        return ""

    if isinstance(payload, str):
        return payload.strip()

    if isinstance(payload, dict):
        for key in ("markdown", "md", "result", "text"):
            value = payload.get(key)
            if isinstance(value, str) and value.strip():
                return value.strip()

    if isinstance(payload, list):
        pieces = [_extract_markdown_from_vl_payload(item) for item in payload]
        pieces = [p for p in pieces if p]
        return "\n\n".join(pieces).strip()

    if hasattr(payload, "save_to_markdown") and callable(payload.save_to_markdown):
        try:
            with tempfile.TemporaryDirectory(prefix="ocr_vl_md_") as tmp_dir:
                saved = False
                for call in (
                    lambda: payload.save_to_markdown(tmp_dir),
                    lambda: payload.save_to_markdown(output_dir=tmp_dir),
                    lambda: payload.save_to_markdown(save_dir=tmp_dir),
                ):
                    try:
                        call()
                        saved = True
                        break
                    except TypeError:
                        continue
                if saved:
                    md_files = sorted(Path(tmp_dir).rglob("*.md"))
                    chunks = [p.read_text(encoding="utf-8", errors="replace") for p in md_files]
                    merged = "\n\n".join(c.strip() for c in chunks if c.strip())
                    if merged:
                        return merged
        except Exception:
            logger.debug("save_to_markdown extraction failed", exc_info=True)

    for attr in ("markdown", "md", "text"):
        if hasattr(payload, attr):
            value = getattr(payload, attr)
            if isinstance(value, str) and value.strip():
                return value.strip()

    return ""


def _ocr_with_paddle_vl(img: Image.Image) -> tuple[str, float | None]:
    engine = _get_paddle_engine()
    if engine is None:
        return "", None

    try:
        rgb = img.convert("RGB")
        arr = np.array(rgb) if _HAS_CV2 else None
        if arr is None:
            return "", None

        inference_calls = [
            lambda: engine.predict(arr),
            lambda: engine.predict([arr]),
            lambda: engine.ocr(arr, cls=True),
        ]
        payload = None
        for call in inference_calls:
            try:
                payload = call()
                break
            except Exception:
                continue

        markdown = _extract_markdown_from_vl_payload(payload)
        if markdown:
            return markdown, 98.0
    except Exception:
        logger.debug("PaddleOCR-VL inference failed on page image", exc_info=True)

    return "", None


def _extract_text_from_paddle_result(result: Any) -> tuple[str, float | None]:
    """Normalize PaddleOCR 3.x OCRResult output and legacy tuple output."""
    if not result:
        return "", None

    entries = result if isinstance(result, list) else [result]
    texts: list[str] = []
    scores: list[float] = []

    for entry in entries:
        rec_texts = None
        rec_scores = None

        if isinstance(entry, dict):
            rec_texts = entry.get("rec_texts")
            rec_scores = entry.get("rec_scores")
        else:
            rec_texts = getattr(entry, "rec_texts", None)
            rec_scores = getattr(entry, "rec_scores", None)

        if isinstance(rec_texts, (list, tuple)):
            texts.extend(str(text).strip() for text in rec_texts if str(text).strip())
        elif isinstance(rec_texts, str) and rec_texts.strip():
            texts.append(rec_texts.strip())

        if isinstance(rec_scores, (list, tuple)):
            for score in rec_scores:
                if score is None:
                    continue
                score_value = float(score)
                scores.append(score_value * 100 if score_value <= 1.0 else score_value)

    if texts:
        avg = round(sum(scores) / len(scores), 2) if scores else None
        return "\n".join(texts).strip(), avg

    legacy_lines: list[str] = []
    legacy_scores: list[float] = []
    if entries and isinstance(entries[0], list):
        for item in entries[0] or []:
            if not item or len(item) < 2:
                continue
            txt = str(item[1][0]).strip()
            conf = float(item[1][1]) if item[1][1] is not None else None
            if txt:
                legacy_lines.append(txt)
            if conf is not None:
                legacy_scores.append(conf * 100 if conf <= 1.0 else conf)

    if legacy_lines:
        avg = round(sum(legacy_scores) / len(legacy_scores), 2) if legacy_scores else None
        return "\n".join(legacy_lines).strip(), avg

    return "", None


def _is_tesseract_available() -> bool:
    global _tesseract_checked, _tesseract_available
    if _tesseract_checked:
        return _tesseract_available
    _tesseract_checked = True
    try:
        pytesseract.get_tesseract_version()
        _tesseract_available = True
    except Exception:
        _tesseract_available = False
        logger.warning("Tesseract not available; OCR will use PaddleOCR/digital text only")
    return _tesseract_available


def _merge_text_digital_first(digital_text: str, ocr_text: str) -> str:
    """Merge two text blocks while preserving digital lines and removing duplicates."""
    ordered_lines: list[str] = []
    seen: set[str] = set()
    for block in (digital_text, ocr_text):
        for line in block.splitlines():
            normalized = re.sub(r"\s+", " ", line.strip().lower())
            if not normalized:
                continue
            if normalized in seen:
                continue
            seen.add(normalized)
            ordered_lines.append(line.strip())
    return "\n".join(ordered_lines)


# ================================================================== pre-processing

def _preprocess(img: Image.Image, aggressive: bool = False) -> Image.Image:
    """Advanced image preprocessing pipeline for OCR accuracy."""
    if not _HAS_CV2:
        # PIL-only fallback: sharpen + contrast
        img = img.convert("L")
        img = img.filter(ImageFilter.SHARPEN)
        enhancer = ImageEnhance.Contrast(img)
        img = enhancer.enhance(1.5)
        return img

    arr = np.array(img.convert("RGB"))
    gray = cv2.cvtColor(arr, cv2.COLOR_RGB2GRAY)

    # Step 1: Denoise
    denoised = cv2.fastNlMeansDenoising(gray, h=10, templateWindowSize=7, searchWindowSize=21)

    # Step 2: CLAHE contrast enhancement
    clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8))
    enhanced = clahe.apply(denoised)

    # Step 3: Adaptive thresholding (handles uneven lighting / shadows)
    if aggressive:
        binary = cv2.adaptiveThreshold(
            enhanced, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
            cv2.THRESH_BINARY, 31, 10,
        )
    else:
        # Otsu's method works well for clean documents
        _, binary = cv2.threshold(enhanced, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)

    # Step 4: Morphological close (fill small gaps in letters)
    kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (2, 2))
    closed = cv2.morphologyEx(binary, cv2.MORPH_CLOSE, kernel)

    # Step 5: Deskew
    deskewed = _deskew(closed)

    return Image.fromarray(deskewed)


def _preprocess_light(img: Image.Image) -> Image.Image:
    """Low-cost preprocessing for fast OCR backends like PaddleOCR."""
    if not _HAS_CV2:
        img = img.convert("L")
        enhancer = ImageEnhance.Contrast(img)
        return enhancer.enhance(1.1)

    arr = np.array(img.convert("RGB"))
    gray = cv2.cvtColor(arr, cv2.COLOR_RGB2GRAY)
    clahe = cv2.createCLAHE(clipLimit=1.5, tileGridSize=(8, 8))
    enhanced = clahe.apply(gray)
    return Image.fromarray(enhanced)


def _deskew(gray: Any) -> Any:
    """Detect skew angle from text lines and rotate to correct it."""
    _, binary_inv = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY_INV + cv2.THRESH_OTSU)

    coords = np.column_stack(np.where(binary_inv > 0))
    if coords.shape[0] < 50:
        return gray

    angle = cv2.minAreaRect(coords)[-1]
    if angle < -45:
        angle = -(90 + angle)
    else:
        angle = -angle

    if abs(angle) < 0.5:
        return gray

    h, w = gray.shape
    center = (w // 2, h // 2)
    mat = cv2.getRotationMatrix2D(center, angle, 1.0)
    rotated = cv2.warpAffine(
        gray, mat, (w, h), flags=cv2.INTER_CUBIC, borderMode=cv2.BORDER_REPLICATE,
    )
    logger.debug("Deskewed by %.2f deg", angle)
    return rotated


def _upscale_if_small(img: Image.Image, min_dpi_equiv: int = 300) -> Image.Image:
    """Upscale very small images so Tesseract gets enough pixel detail."""
    w, h = img.size
    if w >= 1500:
        return img
    if w < 600 or h < 600:
        scale = max(2, min_dpi_equiv // min(w, h) + 1)
        img = img.resize((w * scale, h * scale), Image.LANCZOS)
        logger.debug("Upscaled small image %dx%d -> %dx%d", w, h, w * scale, h * scale)
    return img


# ================================================================== worker startup warmup

def prewarm_ocr_engines() -> None:
    """Pre-initialize OCR engines on worker startup to avoid runtime delays.
    
    Called once per worker process to amortize model loading cost.
    This ensures PaddleOCR is loaded into memory early, so individual image
    processing doesn't incur initialization overhead.
    """
    logger.info("[OCR] Prewarming OCR engines on worker startup...")
    
    # Pre-warm PaddleOCR
    if settings.enable_paddle_ocr:
        try:
            engine = _get_paddle_engine()
            if engine:
                logger.info("[OCR] Pre-warmed PaddleOCR engine")
        except Exception:
            logger.warning("[OCR] Failed to pre-warm PaddleOCR", exc_info=True)
    
    # Keep EasyOCR cold unless PaddleOCR fails and we truly need the fallback.
    if settings.easyocr_enabled:
        logger.info("[OCR] EasyOCR fallback left cold (will load only if PaddleOCR returns no text)")
    
    logger.info("[OCR] OCR engines prewarmed — ready to process images")


# ================================================================== extraction router

def _detect_extractor_for_unknown(path: Path) -> str:
    """Best-effort detection for files with unknown/missing suffix.

    Reads a magic-number header and falls back to extension-less heuristics,
    returning one of: ``"pdf" | "image" | "docx" | "excel" | "text"``. This
    keeps the pipeline alive for files coming from sources that strip or
    mangle extensions (e.g. some scanners, multipart proxies, mobile uploads).
    """
    try:
        with open(path, "rb") as fh:
            header = fh.read(8)
    except OSError:
        return "text"

    if header.startswith(b"%PDF"):
        return "pdf"
    # Common image magic numbers
    if (
        header.startswith(b"\x89PNG")
        or header.startswith(b"\xff\xd8\xff")              # JPEG
        or header[:6] in (b"GIF87a", b"GIF89a")
        or header.startswith(b"BM")                         # BMP
        or header.startswith(b"II*\x00") or header.startswith(b"MM\x00*")  # TIFF
        or header[:4] == b"RIFF"                            # WebP container
    ):
        return "image"
    # Office Open XML formats are ZIP-based: PK\x03\x04
    if header[:2] == b"PK":
        # Try DOCX first; if that fails caller will fall back to excel/text
        return "docx"
    return "text"


def extract_text(file_path: str | Path) -> list[PageResult]:
    """Run extraction on any supported file and return per-page results."""
    path = Path(file_path)
    suffix = path.suffix.lower()
    if suffix in _PDF_EXTENSIONS:
        raw = _extract_from_pdf(path)
    elif suffix in _IMAGE_EXTENSIONS:
        raw = _extract_from_image(path)
    elif suffix in _DOCX_EXTENSIONS:
        raw = _extract_from_docx(path)
    elif suffix in _EXCEL_EXTENSIONS:
        raw = _extract_from_excel(path)
    elif suffix in _PPTX_EXTENSIONS:
        raw = _extract_from_pptx(path)
    elif suffix in _ODT_EXTENSIONS:
        raw = _extract_from_odt(path)
    elif suffix in _RTF_EXTENSIONS:
        raw = _extract_from_rtf(path)
    elif suffix in _TEXT_EXTENSIONS:
        raw = _extract_from_text(path)
    else:
        # Unknown / missing extension: detect by content and try in order.
        kind = _detect_extractor_for_unknown(path)
        raw = None
        attempts = {
            "pdf": _extract_from_pdf,
            "image": _extract_from_image,
            "docx": _extract_from_docx,
            "excel": _extract_from_excel,
            "text": _extract_from_text,
        }
        # Try the detected kind first, then fall back through the rest.
        order = [kind] + [k for k in ("image", "pdf", "docx", "excel", "text") if k != kind]
        last_err: Exception | None = None
        for k in order:
            try:
                raw = attempts[k](path)
                break
            except Exception as exc:
                last_err = exc
                continue
        if raw is None:
            raise ValueError(
                f"Unsupported file type: {suffix or '(no extension)'} — last error: {last_err}"
            )

    # Maintain backward compatibility: callers expect list of (page_num, text, confidence)
    if raw and isinstance(raw[0], dict):
        return [(r.get('page', idx + 1), r.get('text', ''), r.get('confidence')) for idx, r in enumerate(raw)]
    return raw


def extract_text_structured(file_path: str | Path) -> list[dict]:
    """Run extraction and return structured per-page dicts with fields and tables."""
    path = Path(file_path)
    suffix = path.suffix.lower()
    if suffix in _PDF_EXTENSIONS:
        return _extract_from_pdf(path)
    if suffix in _IMAGE_EXTENSIONS:
        return _extract_from_image(path)
    if suffix in _DOCX_EXTENSIONS:
        return _extract_from_docx(path)
    if suffix in _EXCEL_EXTENSIONS:
        return _extract_from_excel(path)
    if suffix in _PPTX_EXTENSIONS:
        return _extract_from_pptx(path)
    if suffix in _ODT_EXTENSIONS:
        return _extract_from_odt(path)
    if suffix in _RTF_EXTENSIONS:
        return _extract_from_rtf(path)
    if suffix in _TEXT_EXTENSIONS:
        return _extract_from_text(path)
    # Unknown / missing extension: detect by content and try in order.
    kind = _detect_extractor_for_unknown(path)
    attempts = {
        "pdf": _extract_from_pdf,
        "image": _extract_from_image,
        "docx": _extract_from_docx,
        "excel": _extract_from_excel,
        "text": _extract_from_text,
    }
    order = [kind] + [k for k in ("image", "pdf", "docx", "excel", "text") if k != kind]
    last_err: Exception | None = None
    for k in order:
        try:
            return attempts[k](path)
        except Exception as exc:
            last_err = exc
            continue
    raise ValueError(
        f"Unsupported file type: {suffix or '(no extension)'} — last error: {last_err}"
    )


# ================================================================== PDF extraction

def _process_pdf_page_worker(pdf_path: str, page_idx: int) -> dict:
    """
    Worker function to process a single PDF page in parallel.
    Returns a dict with page results that can be merged with other pages.
    """
    with pdfplumber.open(pdf_path) as pdf:
        page = pdf.pages[page_idx]
        page_num = page_idx + 1
        
        parts: list[str] = []
        tables_found = []

        # 1. Embedded text
        text = page.extract_text() or ""
        if text.strip():
            parts.append(text.strip())

        # 2. Table extraction (pdfplumber structured tables)
        try:
            tables = page.extract_tables()
            for table in (tables or []):
                table_text = _format_table(table)
                if table_text:
                    parts.append(table_text)
                    tables_found.append(table)
        except Exception:
            logger.debug("Table extraction failed on page %d", page_num)

        digital_text = "\n\n".join(parts).strip()

        # Some scanned PDFs expose tiny or garbled text layers that pdfplumber can read
        # but that still need OCR to get usable content.
        digital_text_len = len(re.sub(r"\s+", "", digital_text))
        should_ocr = settings.enable_secondary_ocr_on_pdf or not digital_text or digital_text_len < 20

        # If pdfplumber already extracted a substantial amount of text or a table,
        # do not pay the full OCR cost unless the page is clearly weak.
        if tables_found and digital_text_len >= 20 and not settings.force_secondary_ocr_on_pdf:
            should_ocr = False

        # 3. OCR pass (for scanned overlays / image-only regions)
        page_text, conf = "", None
        if should_ocr:
            page_text, conf = _ocr_pdf_page(page)

        if digital_text and settings.enable_secondary_ocr_on_pdf:
            merged = _merge_text_digital_first(digital_text, page_text)
            text_for_fields = merged
            confidence = conf if conf is not None else 99.0
        elif digital_text and digital_text_len >= 20:
            text_for_fields = digital_text
            confidence = 99.0
        elif digital_text:
            text_for_fields = page_text or digital_text
            confidence = conf if page_text else 99.0
        else:
            text_for_fields = page_text
            confidence = conf

        parsed = _extract_fields_and_tables(text_for_fields)
        # Merge in tables found by pdfplumber
        if tables_found:
            parsed['tables'] = tables_found + parsed['tables']
        
        return {
            'page': page_num,
            'text': text_for_fields,
            'fields': parsed['fields'],
            'tables': parsed['tables'],
            'confidence': confidence
        }


def _extract_from_pdf(path: Path) -> list[PageResult]:
    """Extract text from PDF with embedded text + table extraction + scanned fallback. Returns list of dicts with text, fields, tables, confidence."""
    # If the PDF is encrypted, decrypt it into a tempfile first using pikepdf
    # (pdfplumber/pdfminer crash with PDFPasswordIncorrect otherwise).
    working_path = _maybe_decrypt_pdf(path)

    # Optional: layout-aware extraction via IBM docling. Returns None when
    # the library isn't installed OR fails on this file, so we fall through
    # to the established pdfplumber+OCR pipeline.
    if getattr(settings, "use_docling", False):
        try:
            from .docling_engine import extract_with_docling
            docling_pages = extract_with_docling(working_path)
            if docling_pages and any(p.get("text") for p in docling_pages):
                logger.info("PDF %s parsed via docling (%d pages)", path.name, len(docling_pages))
                return docling_pages
        except Exception:
            logger.debug("docling path failed, falling back to pdfplumber+OCR", exc_info=True)

    with pdfplumber.open(working_path) as pdf:
        num_pages = len(pdf.pages)
    
    # Process pages in parallel (2-4 workers based on document size)
    max_workers = min(4, max(2, num_pages // 2))
    page_indices = list(range(num_pages))
    
    results: list[dict] = []
    try:
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            page_results = executor.map(
                _process_pdf_page_worker,
                [str(working_path)] * num_pages,
                page_indices
            )
            results = list(page_results)
    except Exception as e:
        logger.error("Parallel PDF page processing failed, falling back to sequential: %s", e)
        # Fallback: process pages sequentially if threading fails
        for page_idx in page_indices:
            try:
                result = _process_pdf_page_worker(str(working_path), page_idx)
                results.append(result)
            except Exception as page_err:
                logger.error("Error processing page %d: %s", page_idx + 1, page_err)
    
    return results


def _maybe_decrypt_pdf(path: Path) -> Path:
    """If ``path`` is an encrypted PDF, decrypt it using pikepdf with a list of\n    candidate passwords and return a path to the decrypted copy. Otherwise\n    return ``path`` unchanged.\n\n    Passwords come from env ``OCR_PDF_PASSWORDS`` (comma-separated). An empty\n    password is always tried first because many "encrypted" PDFs only have\n    permission flags set with no owner password.\n    """
    if not _HAS_PIKEPDF:
        return path

    try:
        with open(path, "rb") as fh:
            head = fh.read(1024)
        if b"/Encrypt" not in head and b"/Encrypt " not in head:
            # Not encrypted
            return path
    except Exception:
        return path

    candidate_passwords: list[str] = [""]
    extra = os.environ.get("OCR_PDF_PASSWORDS", "") or getattr(settings, "pdf_passwords", "")
    for pw in (extra or "").split(","):
        pw = pw.strip()
        if pw and pw not in candidate_passwords:
            candidate_passwords.append(pw)

    last_err: Exception | None = None
    for pw in candidate_passwords:
        try:
            with _pikepdf.open(str(path), password=pw) as pdf:  # type: ignore[union-attr]
                tmp = tempfile.NamedTemporaryFile(
                    delete=False, suffix=".pdf", prefix="ocr_decrypted_",
                )
                tmp.close()
                pdf.save(tmp.name)
                logger.info(
                    "Decrypted encrypted PDF (password=%s) -> %s",
                    "<empty>" if not pw else "<provided>", tmp.name,
                )
                return Path(tmp.name)
        except Exception as exc:  # pragma: no cover - depends on pikepdf version
            last_err = exc
            continue

    logger.warning(
        "Could not decrypt PDF %s with any candidate password (set OCR_PDF_PASSWORDS); "
        "downstream parser may fail. Last error: %s",
        path, last_err,
    )
    return path


def _extract_from_pdf_old(path: Path) -> list[PageResult]:
    """Extract text from PDF with embedded text + table extraction + scanned fallback. Returns list of dicts with text, fields, tables, confidence."""
    results: list[dict] = []

    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            parts: list[str] = []
            tables_found = []

            # 1. Embedded text
            text = page.extract_text() or ""
            if text.strip():
                parts.append(text.strip())

            # 2. Table extraction (pdfplumber structured tables)
            try:
                tables = page.extract_tables()
                for table in (tables or []):
                    table_text = _format_table(table)
                    if table_text:
                        parts.append(table_text)
                        tables_found.append(table)
            except Exception:
                logger.debug("Table extraction failed on page %d", i)

            digital_text = "\n\n".join(parts).strip()

            # Some scanned PDFs expose tiny or garbled text layers that pdfplumber can read
            # but that still need OCR to get usable content.
            digital_text_len = len(re.sub(r"\s+", "", digital_text))
            should_ocr = settings.enable_secondary_ocr_on_pdf or not digital_text or digital_text_len < 20

            # If pdfplumber already extracted a substantial amount of text or a table,
            # do not pay the full OCR cost unless the page is clearly weak.
            if tables_found and digital_text_len >= 20 and not settings.force_secondary_ocr_on_pdf:
                should_ocr = False

            # 3. OCR pass (for scanned overlays / image-only regions)
            page_text, conf = "", None
            if should_ocr:
                page_text, conf = _ocr_pdf_page(page)

            if digital_text and settings.enable_secondary_ocr_on_pdf:
                merged = _merge_text_digital_first(digital_text, page_text)
                text_for_fields = merged
                confidence = conf if conf is not None else 99.0
            elif digital_text and digital_text_len >= 20:
                text_for_fields = digital_text
                confidence = 99.0
            elif digital_text:
                text_for_fields = page_text or digital_text
                confidence = conf if page_text else 99.0
            else:
                text_for_fields = page_text
                confidence = conf

            parsed = _extract_fields_and_tables(text_for_fields)
            # Merge in tables found by pdfplumber
            if tables_found:
                parsed['tables'] = tables_found + parsed['tables']
            results.append({'page': i, 'text': text_for_fields, 'fields': parsed['fields'], 'tables': parsed['tables'], 'confidence': confidence})

    return results


def _ocr_pdf_page(page) -> tuple[str, float | None]:
    """OCR a single PDF page by rendering to image."""
    img = page.to_image(resolution=200).original
    img = _upscale_if_small(img)

    # Try EasyOCR first as primary engine
    _ensure_easyocr_reader()
    if _HAS_EASYOCR and _easyocr_reader is not None:
        try:
            import numpy as np
            arr = np.array(img.convert("RGB"))
            result = _easyocr_reader.readtext(arr, detail=0, paragraph=True)
            easy_text = "\n".join(result).strip()
            if easy_text:
                return easy_text, None
        except Exception:
            logger.debug("EasyOCR inference failed on PDF page image", exc_info=True)

    # PaddleOCR as fallback when EasyOCR did not help
    paddle_text, paddle_conf = _ocr_with_paddle(_preprocess_light(img))
    if paddle_text.strip():
        return paddle_text, paddle_conf

    if not _is_tesseract_available():
        return "", None

    # First pass: standard preprocessing
    cleaned = _preprocess(img, aggressive=False)
    try:
        data = pytesseract.image_to_data(cleaned, output_type=pytesseract.Output.DICT)
    except TesseractNotFoundError:
        logger.warning("Tesseract not found during PDF OCR; skipping fallback OCR for this page")
        return "", None
    text, conf = _aggregate_tesseract_data(data)

    # If confidence is low, retry with aggressive preprocessing
    if conf is not None and conf < 60:
        logger.debug("Low confidence (%.1f) — retrying with aggressive preprocessing", conf)
        cleaned_agg = _preprocess(img, aggressive=True)
        try:
            data2 = pytesseract.image_to_data(cleaned_agg, output_type=pytesseract.Output.DICT)
        except TesseractNotFoundError:
            return text, conf
        text2, conf2 = _aggregate_tesseract_data(data2)
        if conf2 is not None and (conf is None or conf2 > conf):
            text, conf = text2, conf2

    return text, conf


def _format_table(table: list) -> str:
    """Format a pdfplumber table (list of rows) into readable text."""
    if not table:
        return ""
    lines: list[str] = []
    for row in table:
        cells = [str(c).strip() if c else "" for c in row]
        lines.append(" | ".join(cells))
    return "\n".join(lines)


# ================================================================== image extraction

def _extract_from_image(path: Path) -> list[PageResult]:
    """Extract text from image with multi-pass OCR for best results. Returns list of dicts with text, fields, tables, confidence."""
    img = Image.open(path)
    img = _upscale_if_small(img)

    # Handle multi-frame images (e.g. multi-page TIFF)
    results: list[dict] = []
    try:
        n_frames = getattr(img, "n_frames", 1)
    except Exception:
        n_frames = 1

    for frame_idx in range(n_frames):
        try:
            img.seek(frame_idx)
        except EOFError:
            break

        frame = img.copy()

        # Try EasyOCR first as primary engine
        _ensure_easyocr_reader()  # ensure reader is initialized (fallback if pre-warming failed)
        if _HAS_EASYOCR and _easyocr_reader is not None:
            import numpy as np
            arr = np.array(frame.convert("RGB"))
            try:
                logger.info("[OCR] EasyOCR primary used for image frame %s", frame_idx + 1)
                result = _easyocr_reader.readtext(arr, detail=0, paragraph=True)
                text = "\n".join(result)
                conf = None  # EasyOCR does not provide confidence by default
                if text.strip():
                    parsed = _extract_fields_and_tables(text)
                    results.append({'page': frame_idx + 1, 'text': text, 'fields': parsed['fields'], 'tables': parsed['tables'], 'confidence': conf})
                    continue
            except Exception:
                logger.debug("EasyOCR inference failed on image frame", exc_info=True)

        # Try PaddleOCR as fallback
        paddle_text, paddle_conf = _ocr_with_paddle(_preprocess_light(frame))
        if paddle_text.strip():
            parsed = _extract_fields_and_tables(paddle_text)
            results.append({'page': frame_idx + 1, 'text': paddle_text, 'fields': parsed['fields'], 'tables': parsed['tables'], 'confidence': paddle_conf})
            continue

        # Final fallback to Tesseract
        if not _is_tesseract_available():
            results.append({'page': frame_idx + 1, 'text': '', 'fields': {}, 'tables': [], 'confidence': None})
            continue

        cleaned = _preprocess(frame, aggressive=False)
        try:
            data = pytesseract.image_to_data(cleaned, output_type=pytesseract.Output.DICT)
        except TesseractNotFoundError:
            results.append({'page': frame_idx + 1, 'text': '', 'fields': {}, 'tables': [], 'confidence': None})
            continue
        text, conf = _aggregate_tesseract_data(data)

        if conf is not None and conf < 60:
            cleaned_agg = _preprocess(frame, aggressive=True)
            try:
                data2 = pytesseract.image_to_data(cleaned_agg, output_type=pytesseract.Output.DICT)
            except TesseractNotFoundError:
                parsed = _extract_fields_and_tables(text)
                results.append({'page': frame_idx + 1, 'text': text, 'fields': parsed['fields'], 'tables': parsed['tables'], 'confidence': conf})
                continue
            text2, conf2 = _aggregate_tesseract_data(data2)
            if conf2 is not None and (conf is None or conf2 > conf):
                text, conf = text2, conf2

        parsed = _extract_fields_and_tables(text)
        results.append({'page': frame_idx + 1, 'text': text, 'fields': parsed['fields'], 'tables': parsed['tables'], 'confidence': conf})

    return results if results else [{'page': 1, 'text': '', 'fields': {}, 'tables': [], 'confidence': None}]


def _ocr_with_paddle(img: Image.Image) -> tuple[str, float | None]:
    engine = _get_paddle_engine()
    if engine is None:
        return "", None

    if _paddle_engine_kind == "vl":
        return _ocr_with_paddle_vl(img)
    try:
        rgb = img.convert("RGB")
        arr = np.array(rgb) if _HAS_CV2 else None
        if arr is None:
            return "", None
        result = engine.predict(
            arr,
            use_doc_orientation_classify=False,
            use_doc_unwarping=False,
            use_textline_orientation=False,
            text_rec_score_thresh=0.0,
        )
        return _extract_text_from_paddle_result(result)
    except Exception:
        logger.debug("PaddleOCR inference failed on page image", exc_info=True)
        return "", None


# ================================================================== module-level initialization
# Pre-warm OCR engines on module import to avoid per-claim latency
if os.environ.get("DISABLE_OCR_PREWARM") != "1":
    try:
        prewarm_ocr_engines()
    except Exception:
        logger.debug("[OCR] Failed to pre-warm on module import (will lazy-load on first use)", exc_info=True)


# ================================================================== DOCX extraction

def _extract_from_docx(path: Path) -> list[PageResult]:
    """Extract text from Word documents including paragraphs and tables."""
    if not _HAS_DOCX:
        raise ValueError("python-docx not installed  -- cannot process .docx files")

    doc = _docx.Document(str(path))
    parts: list[str] = []

    # Extract all paragraphs
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)

    # Extract tables
    for table in doc.tables:
        rows: list[str] = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            parts.append("\n".join(rows))

    combined = "\n\n".join(parts)
    return [(1, combined, 99.0)]


# ================================================================== Excel extraction

def _extract_from_excel(path: Path) -> list[PageResult]:
    """Extract text from all sheets in an Excel workbook."""
    if not _HAS_OPENPYXL:
        raise ValueError("openpyxl not installed -- cannot process .xlsx files")

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    results: list[PageResult] = []

    for sheet_idx, sheet_name in enumerate(wb.sheetnames, start=1):
        ws = wb[sheet_name]
        lines: list[str] = []
        lines.append(f"[Sheet: {sheet_name}]")

        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() if c is not None else "" for c in row]
            if any(cells):
                lines.append(" | ".join(cells))

        text = "\n".join(lines)
        results.append((sheet_idx, text, 99.0))

    wb.close()
    return results if results else [(1, "", None)]


# ================================================================== PowerPoint extraction

def _extract_from_pptx(path: Path) -> list[PageResult]:
    """Extract text from PowerPoint slides (one entry per slide)."""
    if not _HAS_PPTX:
        raise ValueError("python-pptx not installed -- cannot process .pptx files")

    prs = _Presentation(str(path))  # type: ignore[misc]
    results: list[PageResult] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            # Plain text frames
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
            # Tables on slides
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [(cell.text or "").strip() for cell in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
        # Speaker notes
        try:
            notes = slide.notes_slide.notes_text_frame.text  # type: ignore[union-attr]
            if notes and notes.strip():
                parts.append(f"[Speaker notes]\n{notes.strip()}")
        except Exception:
            pass
        text = "\n".join(parts).strip()
        results.append((slide_idx, text, 99.0))
    return results if results else [(1, "", None)]


# ================================================================== OpenDocument extraction

def _extract_from_odt(path: Path) -> list[PageResult]:
    """Extract text from OpenDocument files (.odt / .ods / .odp)."""
    if not _HAS_ODF:
        raise ValueError("odfpy not installed -- cannot process OpenDocument files")

    doc = _odf_load(str(path))  # type: ignore[misc]
    parts: list[str] = []

    # Paragraphs (text + presentation slides)
    for p in doc.getElementsByType(_odf_text.P):  # type: ignore[attr-defined]
        line = _odf_teletype.extractText(p).strip()  # type: ignore[attr-defined]
        if line:
            parts.append(line)

    # Tables (spreadsheets + tables in text docs)
    for tbl in doc.getElementsByType(_odf_table.Table):  # type: ignore[attr-defined]
        for row in tbl.getElementsByType(_odf_table.TableRow):  # type: ignore[attr-defined]
            cells = []
            for cell in row.getElementsByType(_odf_table.TableCell):  # type: ignore[attr-defined]
                cells.append(_odf_teletype.extractText(cell).strip())  # type: ignore[attr-defined]
            if any(cells):
                parts.append(" | ".join(cells))

    text = "\n".join(parts).strip()
    return [(1, text, 99.0)]


# ================================================================== RTF extraction

def _extract_from_rtf(path: Path) -> list[PageResult]:
    """Strip RTF control codes and return clean text."""
    raw = path.read_text(encoding="utf-8", errors="replace")
    if _HAS_STRIPRTF and _rtf_to_text is not None:
        try:
            text = _rtf_to_text(raw, errors="ignore").strip()
        except Exception:
            logger.debug("striprtf parse failed, returning raw payload", exc_info=True)
            text = raw
    else:
        # Best-effort fallback: drop \controlwords and braces. Better than serving raw RTF.
        text = re.sub(r"\\[a-zA-Z]+-?\d* ?", " ", raw)
        text = text.replace("{", " ").replace("}", " ")
        text = re.sub(r"\s+", " ", text).strip()
    return [(1, text, 99.0)]


# ================================================================== text / CSV / JSON extraction

def _extract_from_text(path: Path) -> list[PageResult]:
    """Read plain text, CSV, JSON, XML, HTML files."""
    suffix = path.suffix.lower()

    # Try utf-8 first, then latin-1 as fallback
    for encoding in ("utf-8", "latin-1"):
        try:
            raw = path.read_text(encoding=encoding)
            break
        except (UnicodeDecodeError, ValueError):
            continue
    else:
        raw = path.read_bytes().decode("utf-8", errors="replace")

    if suffix == ".csv":
        return _extract_from_csv_text(raw)
    if suffix == ".json":
        return _extract_from_json_text(raw)

    # Plain text / XML / HTML / MD — return as-is
    return [(1, raw.strip(), 99.0)]


def _extract_from_csv_text(raw: str) -> list[PageResult]:
    """Parse CSV into readable tabular text."""
    reader = csv.reader(io.StringIO(raw))
    lines: list[str] = []
    for row in reader:
        cells = [c.strip() for c in row]
        if any(cells):
            lines.append(" | ".join(cells))
    return [(1, "\n".join(lines), 99.0)]


def _extract_from_json_text(raw: str) -> list[PageResult]:
    """Flatten JSON into readable text."""
    try:
        data = json.loads(raw)
        text = json.dumps(data, indent=2, ensure_ascii=False)
    except json.JSONDecodeError:
        text = raw
    return [(1, text, 99.0)]


# ================================================================== Tesseract helpers

def _aggregate_tesseract_data(data: dict) -> tuple[str, float | None]:
    """Combine Tesseract word-level output into full text + average confidence."""
    words: list[str] = []
    confidences: list[float] = []
    for txt, c in zip(data["text"], data["conf"], strict=False):
        c = float(c)
        if c < 0:
            continue
        stripped = txt.strip()
        if stripped:
            words.append(stripped)
            confidences.append(c)

    text = " ".join(words)
    avg_conf = round(sum(confidences) / len(confidences), 2) if confidences else None
    return text, avg_conf
